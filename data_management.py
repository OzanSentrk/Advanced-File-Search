# data_management.py

import os
import pickle
import pdfplumber
import faiss
import numpy as np
from sklearn.feature_extraction.text import TfidfVectorizer
from sentence_transformers import SentenceTransformer
from constants import db_directory, turkish_stopwords
from preprocessing import preprocess_text
from docx import Document
import pandas as pd
from pptx import Presentation
import time
from pathlib import Path
# Global değişkenler
filenames = []
filename_to_id = {}  # Dosya isimlerine karşılık gelen ID'ler
preprocessed_texts = []
tfidf_vectorizer = None
tfidf_matrix = None
index = None  # FAISS indeksi
dimension = None  # Vektör boyutu

# Windows'ta AppData/Local dizinini alın
INDEX_BASE_DIR = Path.home() / '.afs_indices'

# PDF dosyalarını okuma fonksiyonu
def read_pdf(file_path):
   
    text = ""
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
    except Exception as e:
        print(f"PDF dosyası okunamadı: {file_path}, hata: {e}")
    return text

# Word dosyalarını okuma fonksiyonu
def read_word(file_path):
   
    text = ""
    try:
        doc = Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        print(f"Word dosyası okunamadı: {file_path}, hata: {e}")
    return text

# Excel dosyalarını okuma fonksiyonu
def read_excel(file_path):
    
    text = ""
    try:
        xls = pd.ExcelFile(file_path)
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            for column in df.columns:
                column_data = ' '.join(df[column].astype(str).tolist())
                text += column_data + "\n"
    except Exception as e:
        print(f"Excel dosyası okunamadı: {file_path}, hata: {e}")
    return text

# PowerPoint dosyalarını okuma fonksiyonu
def read_powerpoint(file_path):
    
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"PowerPoint dosyası okunamadı: {file_path}, hata: {e}")
    return text

# FAISS indeksini IDMap ile oluşturma fonksiyonu
def create_faiss_index(dimension):
    # IndexFlatIP ile bir indeks oluşturun
    flat_index = faiss.IndexFlatIP(dimension)
    # IDMap2 ile indeksinizi sarmalayın
    id_map = faiss.IndexIDMap2(flat_index)
    return id_map

# Dosyaların içeriğini alma fonksiyonu
def get_file_contents(folder):
    
    files_content = []
    filenames_local = []

    # Hedef klasörün adını güvenli bir şekilde alalım
    folder_name = os.path.basename(os.path.normpath(folder))
    # İndeks dosyalarının saklandığı dizini belirleyin
    db_directory = INDEX_BASE_DIR / folder_name

    for root, dirs, files in os.walk(folder):
        # İndeks dizinini atlayın
        dirs[:] = [d for d in dirs if os.path.join(root, d) != str(db_directory)]
        for file in files:
            file_path = os.path.join(root, file)
            
            content_raw = ""
            if file.lower().endswith('.txt'):
                try:
                    with open(file_path, 'r', encoding='utf-8') as f:
                        content_raw = f.read()
                        
                except Exception as e:
                        print(f"Metin dosyası okunamadı: {file_path}, hata: {e}")
            elif file.lower().endswith('.pdf'):
                    content_raw = read_pdf(file_path)
            elif file.lower().endswith(('.docx', '.doc')):
                    content_raw = read_word(file_path)
            elif file.lower().endswith(('.xlsx', '.xls')):
                    content_raw = read_excel(file_path)
            elif file.lower().endswith(('.pptx', '.ppt')):
                    content_raw = read_powerpoint(file_path)
            else:
                    
                    continue  # Desteklenmeyen dosyaları atla
                # Metin ön işlemesi
            content = preprocess_text(content_raw)
            if content:
                    files_content.append(content)
                    # Dosya yolunu hedef klasöre göre göreceli hale getirin
                    file_rel_path = os.path.relpath(file_path, folder)
                    filenames_local.append(file_rel_path)
    
    return files_content, filenames_local


# Vektörleştirme ve indeksleme fonksiyonu
def vectorize_and_index_content(contents):
    global index, tfidf_vectorizer, tfidf_matrix, dimension, filenames, preprocessed_texts, filename_to_id
   
    
    # TF-IDF vektörleştirme
    tfidf_vectorizer = TfidfVectorizer(stop_words=list(turkish_stopwords))
    tfidf_matrix = tfidf_vectorizer.fit_transform(contents)
   
    
    # TF-IDF matrisini yoğunlaştırma ve float32 tipine dönüştürme
    tfidf_matrix_dense = tfidf_matrix.toarray().astype('float32')
    faiss.normalize_L2(tfidf_matrix_dense)
    
    # Embedding boyutunu belirle
    dimension = tfidf_matrix_dense.shape[1]
    
    # FAISS indeksini IDMap ile oluştur
    index = create_faiss_index(dimension)
    
    # FAISS indeksine ekle
    # ID'ler için mevcut filename_to_id'den maksimum ID'yi bul ve artan ID'ler ata
    if filename_to_id:
        max_id = max(filename_to_id.values())
    else:
        max_id = -1  # İlk ID 0 olacak şekilde başlat
    
    ids = np.arange(max_id + 1, max_id + 1 + len(filenames)).astype('int64')
    index.add_with_ids(tfidf_matrix_dense, ids)
    
    # filename_to_id'yi güncelle
    for filename, id_ in zip(filenames, ids):
        filename_to_id[filename] = id_
    
   
    
    return tfidf_matrix

# İndeksi ve dosya isimlerini kaydetme fonksiyonu
def save_index_and_filenames(filenames_local, preprocessed_texts_local, folder):
    global index, tfidf_vectorizer, filename_to_id

    # Hedef klasörün adını güvenli bir şekilde alın
    folder_name = os.path.basename(os.path.normpath(folder))
    # İndeks dosyalarını saklamak için klasör oluşturun
    db_directory = INDEX_BASE_DIR / folder_name
    db_directory.mkdir(parents=True, exist_ok=True)

    # FAISS indeks yolunu ayarlayın
    faiss_index_path = db_directory / 'faiss_index.index'
   

    # FAISS indeksini kaydet
    faiss.write_index(index, str(faiss_index_path))

    # Diğer verileri kaydet
    with open(db_directory / 'folder_path.txt', 'w', encoding='utf-8') as f:
        f.write(folder)
    
    with open(db_directory / 'filenames.pkl', 'wb') as f:
        pickle.dump(filenames_local, f)
    with open(db_directory / 'filename_to_id.pkl', 'wb') as f:
        pickle.dump(filename_to_id, f)
    with open(db_directory / 'preprocessed_texts.pkl', 'wb') as f:
        pickle.dump(preprocessed_texts_local, f)
    with open(db_directory / 'tfidf_vectorizer.pkl', 'wb') as f:
        pickle.dump(tfidf_vectorizer, f)




# İndeksi ve dosya isimlerini yükleme fonksiyonu
def load_index_and_filenames(folder):
    global index, tfidf_vectorizer, filenames, preprocessed_texts, filename_to_id

    # Hedef klasörün adını güvenli bir şekilde alın
    folder_name = os.path.basename(os.path.normpath(folder))
    # İndeks dosyalarının saklandığı dizini belirleyin
    db_directory = INDEX_BASE_DIR / folder_name

    faiss_index_path = db_directory / 'faiss_index.index'
    filenames_path = db_directory / 'filenames.pkl'
    filename_to_id_path = db_directory / 'filename_to_id.pkl'
    preprocessed_texts_path = db_directory / 'preprocessed_texts.pkl'
    tfidf_vectorizer_path = db_directory / 'tfidf_vectorizer.pkl'

    try:
        if all(path.exists() for path in [faiss_index_path, filenames_path, filename_to_id_path, preprocessed_texts_path, tfidf_vectorizer_path]):
            # FAISS indeksini yükleyin
            index = faiss.read_index(str(faiss_index_path))
           

            # Dosya isimlerini yükleyin
            with open(filenames_path, 'rb') as f:
                loaded_filenames = pickle.load(f)
                filenames.clear()
                filenames.extend(loaded_filenames)

            # filename_to_id'yi yükleyin
            with open(filename_to_id_path, 'rb') as f:
                loaded_filename_to_id = pickle.load(f)
                filename_to_id.clear()
                filename_to_id.update(loaded_filename_to_id)

            # Preprocessed metinleri yükleyin
            with open(preprocessed_texts_path, 'rb') as f:
                loaded_preprocessed_texts = pickle.load(f)
                preprocessed_texts.clear()
                preprocessed_texts.extend(loaded_preprocessed_texts)

            # TF-IDF vektörleştiriciyi yükleyin
            with open(tfidf_vectorizer_path, 'rb') as f:
                loaded_tfidf_vectorizer = pickle.load(f)
                tfidf_vectorizer = loaded_tfidf_vectorizer

          
            return False  # İndeksleme yapılmadı, mevcut veri yüklendi
        else:
            
            contents, filenames_local = get_file_contents(folder)
            preprocessed_texts_local = contents
            filenames.clear()
            filenames.extend(filenames_local)
            preprocessed_texts.clear()
            preprocessed_texts.extend(preprocessed_texts_local)
            tfidf_matrix = vectorize_and_index_content(contents)
            save_index_and_filenames(filenames, preprocessed_texts_local, folder)
            return True  # İndeksleme yapıldı
    except Exception as e:
        
        raise
# Yeniden Tara fonksiyonu
def rescan_folder(folder):
    global index, tfidf_vectorizer, filenames, preprocessed_texts, filename_to_id

    # Hedef klasörün adını güvenli bir şekilde alalım
    folder_name = os.path.basename(os.path.normpath(folder))
    # İndeks dosyalarının saklandığı dizini belirleyin
    db_directory = INDEX_BASE_DIR / folder_name

   

    # Mevcut dosyaları listele
    current_files = set()
    for root, dirs, files in os.walk(folder):
        # İndeks dizinini atlayın
        dirs[:] = [d for d in dirs if os.path.join(root, d) != str(db_directory)]
        for file in files:
            current_files.add(os.path.relpath(os.path.join(root, file), folder))

    # Silinen dosyaları belirle
    deleted_files = set(filenames) - current_files
    # Eklenen dosyaları belirle
    added_files = current_files - set(filenames)

   

    # Silinen dosyaları indeksden çıkar
    if deleted_files:
        ids_to_remove = []
        for file in deleted_files:
            try:
                faiss_id = filename_to_id[file]
                ids_to_remove.append(faiss_id)
                
            except KeyError:
                print(f"Dosya ID'si bulunamadı: {file}")

        if ids_to_remove:
            # FAISS indeksinden sil
            faiss_ids = np.array(ids_to_remove).astype('int64')
            index.remove_ids(faiss.IDSelectorBatch(faiss_ids))
            

        # `filenames`, `preprocessed_texts`, ve `filename_to_id` listelerinden silinen dosyaları kaldır
        filenames = [f for f in filenames if f not in deleted_files]
        preprocessed_texts = [t for f, t in zip(filenames, preprocessed_texts) if f not in deleted_files]
        for file in deleted_files:
            filename_to_id.pop(file, None)

    # Eklenen dosyaları işle ve indekse ekle
    if added_files:
        new_contents = []
        new_filenames = []
        for file in added_files:
            file_path = os.path.join(folder, file)
            content_raw = ""
            if file.lower().endswith('.txt'):
                try:
                    with open(file_path, 'r', encoding='utf-8') as f:
                        content_raw = f.read()
                        
                except Exception as e:
                   
                    continue
            elif file.lower().endswith('.pdf'):
                content_raw = read_pdf(file_path)
               
            elif file.lower().endswith(('.docx', '.doc')):
                content_raw = read_word(file_path)
               
            elif file.lower().endswith(('.xlsx', '.xls')):
                content_raw = read_excel(file_path)
               
            elif file.lower().endswith(('.pptx', '.ppt')):
                content_raw = read_powerpoint(file_path)
                
            else:
               
                continue  # Desteklenmeyen dosyaları atla

            # Metin ön işlemesi
            content = preprocess_text(content_raw)
            if content:
                new_contents.append(content)
                new_filenames.append(file)

        if new_contents:
            # Yeni dosyaları listeye ekle
            preprocessed_texts.extend(new_contents)
            filenames.extend(new_filenames)

            # TF-IDF vektörleştirme
            tfidf_matrix_new = tfidf_vectorizer.transform(new_contents)
            embeddings_new = tfidf_matrix_new.toarray().astype('float32')
            faiss.normalize_L2(embeddings_new)

            # Yeni dosyaların FAISS ID'lerini belirle
            if filename_to_id:
                max_id = max(filename_to_id.values())
            else:
                max_id = -1  # İlk ID 0 olacak şekilde başlat
            new_ids = np.arange(max_id + 1, max_id + 1 + len(new_filenames)).astype('int64')

            # FAISS indeksine ekle
            index.add_with_ids(embeddings_new, new_ids)
           

            # filename_to_id'yi güncelle
            for filename, id_ in zip(new_filenames, new_ids):
                filename_to_id[filename] = id_

    # İndeksi ve listeleri kaydet
    save_index_and_filenames(filenames, preprocessed_texts, folder)

    

# Sorgu arama fonksiyonu
def search_query(query, top_k=5):
    global index, tfidf_vectorizer, filenames, filename_to_id
    if index is None or tfidf_vectorizer is None:
        
        return [], 0
    # Sorguyu ön işleyin
    query_processed = preprocess_text(query)
    # Sorgunun TF-IDF vektörünü oluşturun
    query_tfidf = tfidf_vectorizer.transform([query_processed])
    # Seyrek matrisi yoğun NumPy array'e dönüştürme ve float32 tipine çevirme
    query_tfidf = query_tfidf.toarray().astype('float32')
    # Vektörü normalleştirme
    faiss.normalize_L2(query_tfidf)
    query_vector = query_tfidf
    # Arama süresini ölçmek için zamanı başlat
    start_time = time.time()
    # FAISS ile arama yapın
    scores, indices = index.search(query_vector, top_k)
    # Arama süresini hesapla
    elapsed_time = time.time() - start_time
    results = []
    for idx, score in zip(indices[0], scores[0]):
        if idx in filename_to_id.values():
            # Dosya isimlerini ID'ye göre bul
            filename = next((f for f, id_ in filename_to_id.items() if id_ == idx), None)
            if filename:
                results.append((filename, score))
    return results, elapsed_time

# Manuel belge ekleme fonksiyonu
def add_manual_documents(file_paths, folder):
    global filenames, preprocessed_texts, tfidf_matrix, index, filename_to_id, tfidf_vectorizer

    # Hedef klasörün adını güvenli bir şekilde alalım
    folder_name = os.path.basename(os.path.normpath(folder))
    # İndeks dosyalarının saklandığı dizini belirleyin
    db_directory = INDEX_BASE_DIR / folder_name

    new_contents = []
    new_filenames = []
    new_ids = []

    for file_path in file_paths:
        # Dosya yolunu hedef klasöre göre göreceli hale getirin
        if os.path.commonpath([folder, file_path]) == folder:
            # Dosya hedef klasörün içindeyse, göreceli yolu kullanın
            file_rel_path = os.path.relpath(file_path, folder)
        else:
            # Dosya hedef klasörün dışındaysa, mutlak yolu kullanın
            file_rel_path = os.path.abspath(file_path)

        if file_rel_path in filenames:
            
            continue  # Dosya zaten eklenmişse atla

        content_raw = ""
        if file_path.lower().endswith('.txt'):
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content_raw = f.read()
                    
            except Exception as e:
                
                continue
        elif file_path.lower().endswith('.pdf'):
            content_raw = read_pdf(file_path)
            
        elif file_path.lower().endswith(('.docx', '.doc')):
            content_raw = read_word(file_path)
           
        elif file_path.lower().endswith(('.xlsx', '.xls')):
            content_raw = read_excel(file_path)
            
        elif file_path.lower().endswith(('.pptx', '.ppt')):
            content_raw = read_powerpoint(file_path)
            
        else:
            
            continue  # Desteklenmeyen dosyaları atla

        # Metin ön işlemesi
        content = preprocess_text(content_raw)
        if content:
            new_contents.append(content)
            new_filenames.append(file_rel_path)

    if new_contents:
        # Yeni dosyaları listeye ekle
        preprocessed_texts.extend(new_contents)
        filenames.extend(new_filenames)

        # TF-IDF vektörleştirme
        tfidf_matrix_new = tfidf_vectorizer.transform(new_contents)
        embeddings_new = tfidf_matrix_new.toarray().astype('float32')
        faiss.normalize_L2(embeddings_new)

        # Yeni dosyaların FAISS ID'lerini belirle
        if filename_to_id:
            max_id = max(filename_to_id.values())
        else:
            max_id = -1  # İlk ID 0 olacak şekilde başlat
        new_ids = np.arange(max_id + 1, max_id + 1 + len(new_filenames)).astype('int64')

        # FAISS indeksine ekle
        index.add_with_ids(embeddings_new, new_ids)
        

        # filename_to_id'yi güncelle
        for filename, id_ in zip(new_filenames, new_ids):
            filename_to_id[filename] = id_

        # İndeksi ve listeleri kaydet
        save_index_and_filenames(filenames, preprocessed_texts, folder)
        
   

def get_filenames():
    return filenames

def reset_globals():
    global filenames, filename_to_id, preprocessed_texts, tfidf_vectorizer, tfidf_matrix, index, dimension
    filenames.clear()
    filename_to_id.clear()
    preprocessed_texts.clear()
    tfidf_vectorizer = None
    tfidf_matrix = None
    index = None
    dimension = None
def get_scanned_folder():
    scanned_folders = []
    if INDEX_BASE_DIR.exists():
        for folder in INDEX_BASE_DIR.iterdir():
            if folder.is_dir():
                folder_name = folder.name
                folder_path = get_original_folder_path(folder_name)
                index_size=get_index_size(folder)
                scanned_folders.append({
                    'name': folder_name,
                    'path': folder_path,
                    'size': index_size
                })
    return scanned_folders
def get_index_size(index_folder):
    total_size = 0
    for file in index_folder.glob('*'):
        if file.is_file():
            total_size += file.stat().st_size
    return total_size

def get_original_folder_path(folder_name):
    folder_index_dir =INDEX_BASE_DIR / folder_name
    folder_path_file = folder_index_dir / 'folder_path.txt'
    if folder_path_file.exists():
        with open(folder_path_file, 'r', encoding = 'utf-8') as f:
            folder_path = f.read().strip()
        return folder_path
    else:
        return None



